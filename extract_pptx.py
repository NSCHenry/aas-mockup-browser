#!/usr/bin/env python3
from pptx import Presentation
import sys

def extract_pptx(filepath):
    prs = Presentation(filepath)

    print(f"=== PRESENTATION: {filepath} ===\n")
    print(f"Total Slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*80}")
        print(f"SLIDE {i+1}")
        print(f"{'='*80}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text)
                print()

            # Extract table data if present
            if shape.has_table:
                table = shape.table
                print("\n[TABLE DATA]")
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    print(" | ".join(row_data))
                print()

if __name__ == "__main__":
    if len(sys.argv) > 1:
        extract_pptx(sys.argv[1])
    else:
        print("Usage: python extract_pptx.py <path_to_pptx>")
